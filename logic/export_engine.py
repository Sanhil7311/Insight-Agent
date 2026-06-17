from __future__ import annotations

import io
import logging
import math
from typing import Any

import pandas as pd
import plotly.express as px
import plotly.io as pio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)

_C_NAVY      = RGBColor(0x1E, 0x27, 0x61)
_C_MIDNIGHT  = RGBColor(0x13, 0x17, 0x3D)
_C_ICE       = RGBColor(0xCA, 0xDC, 0xFC)
_C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_C_CHARCOAL  = RGBColor(0x36, 0x45, 0x4F)
_C_SILVER    = RGBColor(0xF2, 0xF4, 0xF8)
_C_ACCENT    = RGBColor(0x4A, 0x90, 0xD9)
_C_MUTED     = RGBColor(0x8A, 0x99, 0xAA)
_C_SUCCESS   = RGBColor(0x27, 0xAE, 0x60)
_C_WARN      = RGBColor(0xF3, 0x9C, 0x12)

_SLIDE_W = Inches(13.33)
_SLIDE_H = Inches(7.5)

_MARGIN_L = Inches(0.55)
_MARGIN_T = Inches(0.45)
_MARGIN_R = Inches(0.55)
_CONTENT_TOP = Inches(1.35)
_CONTENT_H   = Inches(5.75)
_CONTENT_W   = Inches(12.23)

_FONT_BODY    = "Calibri"
_FONT_HEADING = "Cambria"


def _scaled_pt(text: str, base_pt: int, max_chars: int) -> int:
    if len(text) <= max_chars:
        return base_pt
    ratio = max_chars / max(len(text), 1)
    scaled = int(base_pt * math.sqrt(ratio))
    return max(scaled, 8)


def _bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _rect(slide, left, top, width, height, fill_rgb: RGBColor | None = None, line_rgb: RGBColor | None = None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _txbox(
    slide,
    text: str,
    left, top, width, height,
    font_name: str = _FONT_BODY,
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = _C_CHARCOAL,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    auto_scale: bool = True,
    max_chars_hint: int = 80,
    word_wrap: bool = True,
) -> None:
    if auto_scale:
        font_size = _scaled_pt(text, font_size, max_chars_hint)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.word_wrap = word_wrap
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _header_band(slide, title: str, subtitle: str = "") -> None:
    _rect(slide, Inches(0), Inches(0), _SLIDE_W, Inches(1.15), fill_rgb=_C_NAVY)
    title_fs = _scaled_pt(title, 26, 70)
    _txbox(
        slide, title,
        Inches(0.55), Inches(0.10), Inches(10.5), Inches(0.72),
        font_name=_FONT_HEADING, font_size=title_fs,
        bold=True, color=_C_WHITE, align=PP_ALIGN.LEFT, auto_scale=False,
    )
    if subtitle:
        _txbox(
            slide, subtitle,
            Inches(0.55), Inches(0.78), Inches(10.5), Inches(0.32),
            font_name=_FONT_BODY, font_size=11,
            bold=False, color=_C_ICE, align=PP_ALIGN.LEFT, auto_scale=False,
        )


def _new_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _render_chart_png(spec: dict, df: pd.DataFrame) -> bytes | None:
    chart_type = str(spec.get("chart_type", "")).lower().strip()
    x_col = spec.get("x_axis") or None
    y_col = spec.get("y_axis") or None
    title = spec.get("title", "")

    if x_col and x_col not in df.columns:
        x_col = None
    if y_col and y_col not in df.columns:
        y_col = None

    _dispatch = {
        "bar": px.bar, "scatter": px.scatter, "histogram": px.histogram,
        "box": px.box, "line": px.line, "violin": px.violin,
        "pie": px.pie, "area": px.area,
    }

    try:
        if chart_type == "heatmap":
            num = df.select_dtypes(include="number")
            if num.shape[1] < 2:
                return None
            fig = px.imshow(num.corr(), text_auto=".2f", title=title, aspect="auto",
                            color_continuous_scale="Blues")
        else:
            fn = _dispatch.get(chart_type)
            if fn is None:
                return None
            kw: dict[str, Any] = {"data_frame": df, "title": title}
            if x_col:
                kw["x"] = x_col
            if y_col:
                kw["y"] = y_col
            fig = fn(**kw)

        fig.update_layout(
            paper_bgcolor="rgba(255,255,255,1)",
            plot_bgcolor="rgba(242,244,248,1)",
            font=dict(family=_FONT_BODY, color="#363636", size=12),
            title_font=dict(family=_FONT_HEADING, size=15, color="#1E2761"),
            margin=dict(l=48, r=24, t=52, b=40),
        )
        return pio.to_image(fig, format="png", width=1100, height=620, scale=2)
    except Exception as exc:
        logger.warning("Chart render failed for %r: %s", title, exc)
        return None


def _add_cover_slide(prs: Presentation, file_name: str) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_MIDNIGHT)

    _rect(slide, Inches(0), Inches(0), Inches(0.18), _SLIDE_H, fill_rgb=_C_ACCENT)

    cx = Inches(1.0)
    _txbox(
        slide, "INSIGHT-AGENT",
        cx, Inches(1.8), Inches(11), Inches(0.72),
        font_name=_FONT_BODY, font_size=13, bold=True, color=_C_ICE,
        align=PP_ALIGN.LEFT, auto_scale=False,
    )
    report_label = "Autonomous Data Science Report"
    _txbox(
        slide, report_label,
        cx, Inches(2.42), Inches(11), Inches(1.1),
        font_name=_FONT_HEADING, font_size=38, bold=True, color=_C_WHITE,
        align=PP_ALIGN.LEFT, auto_scale=True, max_chars_hint=35,
    )
    fn_fs = _scaled_pt(file_name, 16, 60)
    _txbox(
        slide, f"Dataset: {file_name}",
        cx, Inches(3.68), Inches(11), Inches(0.45),
        font_name=_FONT_BODY, font_size=fn_fs, bold=False, color=_C_ICE,
        align=PP_ALIGN.LEFT, auto_scale=False,
    )
    _txbox(
        slide, "Powered by CrewAI  ·  Local LLM  ·  AutoML",
        cx, Inches(6.5), Inches(11), Inches(0.45),
        font_name=_FONT_BODY, font_size=11, bold=False, color=_C_MUTED,
        align=PP_ALIGN.LEFT, auto_scale=False,
    )


def _add_kpi_slide(prs: Presentation, analytics: dict, predictive_metrics: dict) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_SILVER)
    _header_band(slide, "At a Glance — Key Metrics", "Summary statistics and model performance")

    shape_val = analytics.get("shape", [0, 0])
    rows = shape_val[0] if len(shape_val) > 0 else 0
    cols_count = shape_val[1] if len(shape_val) > 1 else 0
    num_summary = analytics.get("numeric_summary", {})
    null_counts = analytics.get("null_counts", {})
    total_nulls = sum(null_counts.values())
    num_cols_count = len(num_summary)

    pm = predictive_metrics or {}
    metrics = pm.get("metrics", {})
    target = pm.get("target_column", "")
    skipped = pm.get("skipped", True) or pm.get("error")

    dataset_kpis = [
        ("Rows", f"{rows:,}"),
        ("Columns", str(cols_count)),
        ("Numeric Cols", str(num_cols_count)),
        ("Missing Values", str(total_nulls)),
    ]
    ml_kpis = []
    if not skipped and metrics:
        for k, v in list(metrics.items())[:4]:
            label = k.upper().replace("_", " ")
            val = f"{v:.4f}" if isinstance(v, float) else str(v)
            ml_kpis.append((label, val))

    card_w = Inches(2.6)
    card_h = Inches(1.6)
    card_top = Inches(1.55)
    gap = Inches(0.18)
    start_l = Inches(0.55)

    for i, (label, value) in enumerate(dataset_kpis):
        cl = start_l + i * (card_w + gap)
        _rect(slide, cl, card_top, card_w, card_h, fill_rgb=_C_WHITE)
        val_fs = _scaled_pt(value, 34, 8)
        _txbox(slide, value, cl + Inches(0.18), card_top + Inches(0.18),
               card_w - Inches(0.36), Inches(0.9),
               font_name=_FONT_HEADING, font_size=val_fs, bold=True,
               color=_C_NAVY, align=PP_ALIGN.LEFT, auto_scale=False)
        _txbox(slide, label, cl + Inches(0.18), card_top + Inches(1.05),
               card_w - Inches(0.36), Inches(0.42),
               font_name=_FONT_BODY, font_size=11, bold=False,
               color=_C_MUTED, align=PP_ALIGN.LEFT, auto_scale=False)

    if ml_kpis:
        ml_top = card_top + card_h + Inches(0.35)
        ml_label_top = ml_top - Inches(0.32)
        _txbox(slide, f"AutoML — Predicting '{target}'",
               start_l, ml_label_top, Inches(10), Inches(0.28),
               font_name=_FONT_HEADING, font_size=13, bold=True,
               color=_C_NAVY, auto_scale=False)
        ml_card_w = Inches(2.85)
        for i, (label, value) in enumerate(ml_kpis):
            cl = start_l + i * (ml_card_w + gap)
            _rect(slide, cl, ml_top, ml_card_w, card_h, fill_rgb=_C_NAVY)
            val_fs = _scaled_pt(value, 32, 8)
            _txbox(slide, value, cl + Inches(0.18), ml_top + Inches(0.15),
                   ml_card_w - Inches(0.36), Inches(0.9),
                   font_name=_FONT_HEADING, font_size=val_fs, bold=True,
                   color=_C_WHITE, align=PP_ALIGN.LEFT, auto_scale=False)
            _txbox(slide, label, cl + Inches(0.18), ml_top + Inches(1.05),
                   ml_card_w - Inches(0.36), Inches(0.42),
                   font_name=_FONT_BODY, font_size=11,
                   color=_C_ICE, align=PP_ALIGN.LEFT, auto_scale=False)

    corr_data = analytics.get("correlations", {})
    if corr_data:
        corr_top = Inches(4.45)
        _txbox(slide, "Top Correlations",
               start_l, corr_top, Inches(12.2), Inches(0.3),
               font_name=_FONT_HEADING, font_size=13, bold=True,
               color=_C_NAVY, auto_scale=False)
        items = sorted(corr_data.items(), key=lambda kv: abs(kv[1]), reverse=True)[:4]
        col_w = Inches(2.9)
        for i, (pair, r) in enumerate(items):
            cl = start_l + i * (col_w + gap)
            clr = _C_SUCCESS if r > 0 else _C_WARN
            _rect(slide, cl, corr_top + Inches(0.35), col_w, Inches(0.95), fill_rgb=_C_WHITE)
            r_fs = _scaled_pt(f"{r:+.3f}", 22, 7)
            _txbox(slide, f"{r:+.3f}", cl + Inches(0.15), corr_top + Inches(0.38),
                   Inches(1.0), Inches(0.6),
                   font_name=_FONT_HEADING, font_size=r_fs, bold=True,
                   color=clr, auto_scale=False)
            pair_fs = _scaled_pt(pair, 10, 28)
            _txbox(slide, pair, cl + Inches(1.05), corr_top + Inches(0.44),
                   col_w - Inches(1.2), Inches(0.55),
                   font_name=_FONT_BODY, font_size=pair_fs,
                   color=_C_CHARCOAL, auto_scale=False, word_wrap=True)


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = "") -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_WHITE)
    _header_band(slide, title, subtitle)

    if not bullets:
        return

    body_top = _CONTENT_TOP
    body_h_each = Inches(0.54)
    gap = Inches(0.12)
    left = _MARGIN_L
    width = _CONTENT_W

    for i, bullet in enumerate(bullets[:8]):
        bt = body_top + i * (body_h_each + gap)
        if bt + body_h_each > _CONTENT_TOP + _CONTENT_H:
            break
        text = bullet.lstrip("-•*› ").strip()
        if not text:
            continue
        _rect(slide, left, bt + Inches(0.14), Inches(0.07), Inches(0.22), fill_rgb=_C_ACCENT)
        fs = _scaled_pt(text, 14, 110)
        _txbox(
            slide, text,
            left + Inches(0.22), bt, width - Inches(0.22), body_h_each,
            font_name=_FONT_BODY, font_size=fs, color=_C_CHARCOAL,
            align=PP_ALIGN.LEFT, auto_scale=False, word_wrap=True,
        )


def _add_split_chart_slide(
    prs: Presentation,
    title: str,
    insight_bullets: list[str],
    img_bytes: bytes,
) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_WHITE)
    _header_band(slide, title)

    text_left = _MARGIN_L
    text_w = Inches(4.4)
    chart_left = text_left + text_w + Inches(0.3)
    chart_w = Inches(7.6)
    body_top = _CONTENT_TOP
    body_h_each = Inches(0.52)
    gap = Inches(0.14)

    for i, bullet in enumerate(insight_bullets[:7]):
        bt = body_top + i * (body_h_each + gap)
        if bt + body_h_each > _CONTENT_TOP + _CONTENT_H - Inches(0.1):
            break
        text = bullet.lstrip("-•*› ").strip()
        if not text:
            continue
        _rect(slide, text_left, bt + Inches(0.14), Inches(0.07), Inches(0.22), fill_rgb=_C_ACCENT)
        fs = _scaled_pt(text, 13, 55)
        _txbox(
            slide, text,
            text_left + Inches(0.22), bt, text_w - Inches(0.3), body_h_each,
            font_name=_FONT_BODY, font_size=fs, color=_C_CHARCOAL,
            align=PP_ALIGN.LEFT, auto_scale=False, word_wrap=True,
        )

    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(
        img_stream, chart_left, _CONTENT_TOP,
        width=chart_w, height=_CONTENT_H,
    )


def _add_chart_only_slide(prs: Presentation, title: str, img_bytes: bytes) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_WHITE)
    _header_band(slide, title)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(
        img_stream,
        _MARGIN_L, _CONTENT_TOP,
        width=_CONTENT_W, height=_CONTENT_H,
    )


def _add_fallback_chart_slide(prs: Presentation, spec: dict) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_WHITE)
    title = spec.get("title", "Chart")
    _header_band(slide, title, "Interactive chart available in the web dashboard")
    details = (
        f"Chart type: {spec.get('chart_type', 'N/A')}\n"
        f"X-axis: {spec.get('x_axis', '—')}\n"
        f"Y-axis: {spec.get('y_axis', '—')}"
    )
    _txbox(
        slide, details,
        _MARGIN_L, _CONTENT_TOP + Inches(1.0), _CONTENT_W, Inches(2.0),
        font_name=_FONT_BODY, font_size=14, color=_C_MUTED,
        auto_scale=False, word_wrap=True,
    )


def _add_conclusion_slide(prs: Presentation, recommendations: list[str]) -> None:
    slide = _new_blank_slide(prs)
    _bg(slide, _C_MIDNIGHT)

    _txbox(
        slide, "Actionable Recommendations",
        _MARGIN_L, Inches(0.45), Inches(10), Inches(0.72),
        font_name=_FONT_HEADING, font_size=30, bold=True,
        color=_C_WHITE, align=PP_ALIGN.LEFT, auto_scale=False,
    )
    _txbox(
        slide, "Key next steps for decision-makers",
        _MARGIN_L, Inches(1.08), Inches(10), Inches(0.34),
        font_name=_FONT_BODY, font_size=12,
        color=_C_ICE, align=PP_ALIGN.LEFT, auto_scale=False,
    )

    card_w = Inches(5.8)
    card_h = Inches(1.35)
    gap_x = Inches(0.32)
    gap_y = Inches(0.22)
    start_top = Inches(1.6)
    positions = [
        (_MARGIN_L, start_top),
        (_MARGIN_L + card_w + gap_x, start_top),
        (_MARGIN_L, start_top + card_h + gap_y),
        (_MARGIN_L + card_w + gap_x, start_top + card_h + gap_y),
        (_MARGIN_L, start_top + 2 * (card_h + gap_y)),
        (_MARGIN_L + card_w + gap_x, start_top + 2 * (card_h + gap_y)),
    ]

    recs = [r.lstrip("-•*›0123456789). ").strip() for r in recommendations if r.strip()][:6]

    for idx, rec in enumerate(recs):
        cl, ct = positions[idx]
        _rect(slide, cl, ct, card_w, card_h, fill_rgb=_C_NAVY)
        num_label = f"0{idx + 1}"
        _txbox(
            slide, num_label,
            cl + Inches(0.18), ct + Inches(0.08), Inches(0.5), Inches(0.38),
            font_name=_FONT_HEADING, font_size=14, bold=True,
            color=_C_ACCENT, auto_scale=False,
        )
        fs = _scaled_pt(rec, 13, 80)
        _txbox(
            slide, rec,
            cl + Inches(0.18), ct + Inches(0.45), card_w - Inches(0.36), Inches(0.78),
            font_name=_FONT_BODY, font_size=fs, color=_C_ICE,
            align=PP_ALIGN.LEFT, auto_scale=False, word_wrap=True,
        )


def _extract_section(report_md: str, heading: str) -> list[str]:
    lines = report_md.splitlines()
    in_section = False
    bullets: list[str] = []
    for line in lines:
        stripped = line.strip()
        if stripped.lstrip("#").strip().lower() == heading.lower():
            in_section = True
            continue
        if in_section:
            if stripped.startswith("#"):
                break
            if stripped.startswith(("-", "*", "•", "›")) or (len(stripped) > 4 and stripped[0].isdigit()):
                bullets.append(stripped)
            elif stripped and not stripped.startswith("#"):
                if len(bullets) < 2 or (bullets and not bullets[-1].endswith(stripped)):
                    bullets.append(stripped)
    return [b for b in bullets if b.strip()][:8]


def _extract_recommendations(report_md: str) -> list[str]:
    candidates: list[str] = []
    for heading in ["Actionable Business Recommendations", "Recommendations"]:
        candidates = _extract_section(report_md, heading)
        if candidates:
            break
    return candidates


def _slide_json_to_content_map(slide_json: list[dict]) -> dict[str, list[str]]:
    return {s.get("title", ""): s.get("bullets", []) for s in slide_json}


def build_pptx_bytes(
    slide_json: list[dict],
    chart_specs: list[dict],
    df: pd.DataFrame | None,
    file_name: str,
    analytics: dict | None = None,
    predictive_metrics: dict | None = None,
    report_markdown: str = "",
) -> bytes:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _add_cover_slide(prs, file_name)

    _add_kpi_slide(prs, analytics or {}, predictive_metrics or {})

    content_map = _slide_json_to_content_map(slide_json)

    SKIP_SECTIONS = {"actionable business recommendations", "recommendations"}

    for slide_def in slide_json:
        section_title = slide_def.get("title", "Section")
        if section_title.lower().strip() in SKIP_SECTIONS:
            continue
        bullets = slide_def.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [b.strip() for b in bullets.splitlines() if b.strip()]
        _add_content_slide(prs, section_title, bullets)

    if chart_specs and df is not None and not df.empty:
        analyst_bullets_pool: list[str] = []
        for sd in slide_json:
            analyst_bullets_pool.extend(sd.get("bullets", []))

        for i, spec in enumerate(chart_specs):
            img_bytes = _render_chart_png(spec, df)
            chart_title = spec.get("title", "Chart")
            insight_start = i * 3
            insight_end = insight_start + 3
            insight_bullets = analyst_bullets_pool[insight_start:insight_end]

            if img_bytes:
                if insight_bullets:
                    _add_split_chart_slide(prs, chart_title, insight_bullets, img_bytes)
                else:
                    _add_chart_only_slide(prs, chart_title, img_bytes)
            else:
                _add_fallback_chart_slide(prs, spec)
    elif chart_specs:
        for spec in chart_specs:
            _add_fallback_chart_slide(prs, spec)

    recs = _extract_recommendations(report_markdown)
    if not recs:
        for sd in slide_json:
            if sd.get("title", "").lower().strip() in SKIP_SECTIONS:
                recs = sd.get("bullets", [])
                break
    if recs:
        _add_conclusion_slide(prs, recs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()